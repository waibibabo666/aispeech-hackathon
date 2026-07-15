"""PPTX parser — extract text from PowerPoint slides using python-pptx.

Handles:
- Text boxes / text frames in each shape
- Tables (extracted as markdown-like text for LLM readability)
- Speaker notes (high-value for schedule info)
- Placeholder type labels (title vs body)

Not handled (no text available):
- SmartArt, charts, embedded images (would need OCR fallback — future work)
"""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import BaseParser

logger = logging.getLogger("pptx_parser")

TABLE = MSO_SHAPE_TYPE.TABLE
GROUP = MSO_SHAPE_TYPE.GROUP
PLACEHOLDER_TYPE_LABELS: dict = {
    0: "正文",
    1: "标题",
    2: "副标题",
    3: "页脚",
    4: "幻灯片编号",
    5: "日期",
}


class PptxParser(BaseParser):
    def supports(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == ".pptx"

    def parse(self, file_path: Path) -> str:
        prs = Presentation(str(file_path))
        pages: list[str] = []

        for idx, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []

            # ── Extract text from shapes ──
            for shape in slide.shapes:
                # Tables → markdown-like grid
                if shape.has_table:
                    parts.append(self._extract_table(shape.table))
                    continue

                # Text frames (boxes, placeholders, etc.)
                if shape.has_text_frame:
                    tf = shape.text_frame
                    label = ""
                    if shape.is_placeholder:
                        ph = shape.placeholder_format
                        pt = PLACEHOLDER_TYPE_LABELS.get(ph.type, "占位符")
                        label = f"[{pt}] "
                    paragraphs = [p.text for p in tf.paragraphs if p.text.strip()]
                    if paragraphs:
                        prefix = label if label else ""
                        parts.append(prefix + "\n".join(paragraphs))

                # Group shapes — recurse one level
                if shape.shape_type == GROUP:
                    for child in shape.shapes:
                        if child.has_text_frame:
                            cp = [p.text for p in child.text_frame.paragraphs if p.text.strip()]
                            if cp:
                                parts.append("\n".join(cp))

            # ── Speaker notes ──
            notes_text = ""
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide
                    notes_paragraphs = [
                        p.text for p in notes.notes_text_frame.paragraphs if p.text.strip()
                    ]
                    if notes_paragraphs:
                        notes_text = "[备注] " + "\n".join(notes_paragraphs)
            except Exception:
                pass  # notes slide may exist but have no text

            if notes_text:
                parts.append(notes_text)

            if parts:
                pages.append(f"[Slide {idx}]\n" + "\n\n".join(parts))

        if not pages:
            return "[此PPT中未检测到文字]"

        return "\n\n---\n\n".join(pages)

    @staticmethod
    def _extract_table(table) -> str:
        """Convert a pptx table to a compact markdown-like text grid."""
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "[表格]\n" + "\n".join(rows)
